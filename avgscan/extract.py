"""Tekstlaag- en metadata-extractie per bestandstype.

Levert een lijst (locatie, tekst) plus een metadata-dict. De 'locatie' is een
mensleesbare aanduiding (pagina/tabblad/dia) zodat een bevinding terug te vinden
is in het bronbestand.
"""
from __future__ import annotations

import os

# --------------------------------------------------------------------------------------
# OCR (optioneel, opt-in). Beeld-zonder-tekstlaag (een ingescande brief/formulier) is voor
# de gewone extractie een plaatje; juist daar zit soms het risico (een BSN op een gescand
# formulier). Met OCR aan worden die pagina's alsnog gelezen. On-prem via rapidocr
# (onnxruntime): geen externe aanroep, dus de persoonsgegevens blijven op de machine.
# OCR is TRAAG (seconden per pagina) en niet bit-identiek reproduceerbaar zoals de
# tekstlaag; daarom standaard uit en per bevinding zichtbaar gelabeld als "(OCR)".
_OCR = {"enabled": False, "dpi": 300, "min_conf": 0.5, "max_pages": 20}
_ocr_engine = None


def configure_ocr(enabled=False, dpi=300, min_confidence=0.5, max_pages_per_doc=20):
    """Zet OCR aan/uit en stel de parameters in (aangeroepen bij het opstarten)."""
    _OCR.update(enabled=bool(enabled), dpi=int(dpi),
                min_conf=float(min_confidence), max_pages=int(max_pages_per_doc))


def _get_ocr_engine():
    global _ocr_engine
    if _ocr_engine is None:
        from rapidocr_onnxruntime import RapidOCR   # zware, optionele dependency: lazy
        _ocr_engine = RapidOCR()
    return _ocr_engine


def _ocr_page(page) -> str:
    """OCR één PDF-pagina (PyMuPDF-page) en geef de herkende tekst terug.

    Alleen tekst boven de confidence-drempel telt mee; rommel met lage zekerheid
    (typisch OCR-ruis op tekeningen) valt zo af.
    """
    import numpy as np

    pix = page.get_pixmap(dpi=_OCR["dpi"])
    arr = np.frombuffer(pix.samples, dtype=np.uint8).reshape(pix.height, pix.width, pix.n)
    if pix.n == 4:
        arr = arr[:, :, :3]            # alpha weg
    elif pix.n == 1:
        arr = np.repeat(arr, 3, axis=2)  # grijswaarde -> 3 kanalen
    res, _ = _get_ocr_engine()(arr)
    if not res:
        return ""
    return " ".join(t for _box, t, conf in res if float(conf) >= _OCR["min_conf"])


def _schoon(s):
    """Verwijder ongeldige Unicode (lone surrogates) uit geëxtraheerde tekst.

    Een kapotte tekstlaag kan tekens opleveren die Python niet naar UTF-8 kan schrijven.
    Onbehandeld sloopt dat de SQLite-schrijfactie (UnicodeEncodeError: surrogates not allowed)
    en daarmee de hele run. Gezien 14-07-2026 op een bekendmaking uit 2021.
    """
    if not isinstance(s, str):
        return s
    return s.encode("utf-8", "replace").decode("utf-8", "replace")


def extract(path: str, ext: str):
    """Return (chunks, metadata). chunks = list[(locatie, tekst)]."""
    chunks, meta = _extract_ruw(path, ext)
    chunks = [(_schoon(loc), _schoon(txt)) for loc, txt in chunks]
    meta = {_schoon(k): _schoon(v) for k, v in meta.items()}
    return chunks, meta


def _extract_ruw(path: str, ext: str):
    ext = ext.lower().lstrip(".")
    try:
        if ext == "pdf":
            return _pdf(path)
        if ext == "docx":
            return _docx(path)
        if ext in ("xlsx", "xlsm"):
            return _xlsx(path)
        if ext == "pptx":
            return _pptx(path)
    except Exception as e:  # corrupte/beveiligde bestanden mogen de crawl niet stoppen
        return [], {"_fout": f"{type(e).__name__}: {e}"}
    return [], {}


def _pdf(path):
    import fitz  # PyMuPDF

    chunks, meta = [], {}
    with fitz.open(path) as doc:
        meta = {k: v for k, v in (doc.metadata or {}).items() if v}
        pages_met_beeld = 0
        ocr_gedaan = 0
        for i, page in enumerate(doc, start=1):
            txt = page.get_text("text")
            if txt.strip():
                chunks.append((f"pagina {i}", txt))
            elif page.get_images():
                pages_met_beeld += 1
                # Beeld zonder tekstlaag: met OCR aan alsnog lezen, gelabeld als (OCR).
                if _OCR["enabled"] and ocr_gedaan < _OCR["max_pages"]:
                    try:
                        otxt = _ocr_page(page)
                    except Exception:
                        otxt = ""       # OCR-fout mag de run niet stoppen
                    if otxt.strip():
                        chunks.append((f"pagina {i} (OCR)", otxt))
                        ocr_gedaan += 1
        if pages_met_beeld:
            meta["_ocr_kandidaat_paginas"] = pages_met_beeld
        if ocr_gedaan:
            meta["_ocr_gelezen_paginas"] = ocr_gedaan
    return chunks, meta


def _docx(path):
    import docx

    d = docx.Document(path)
    delen = [p.text for p in d.paragraphs if p.text.strip()]
    for tbl in d.tables:
        for row in tbl.rows:
            cellen = [c.text for c in row.cells if c.text.strip()]
            if cellen:
                delen.append(" | ".join(cellen))
    cp = d.core_properties
    meta = {k: v for k, v in {
        "author": cp.author, "last_modified_by": cp.last_modified_by,
        "title": cp.title, "subject": cp.subject, "comments": cp.comments,
    }.items() if v}
    return [("document", "\n".join(delen))], meta


def _xlsx(path):
    import openpyxl

    wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
    chunks, verborgen = [], []
    for ws in wb.worksheets:
        if ws.sheet_state != "visible":
            verborgen.append(ws.title)  # verborgen/very hidden tabblad = risico
        regels = []
        for row in ws.iter_rows(values_only=True):
            waarden = [str(c) for c in row if c not in (None, "")]
            if waarden:
                regels.append(" | ".join(waarden))
        if regels:
            label = f"tabblad {ws.title}" + (" [VERBORGEN]" if ws.sheet_state != "visible" else "")
            chunks.append((label, "\n".join(regels)))
    props = wb.properties
    meta = {k: v for k, v in {
        "creator": props.creator, "lastModifiedBy": props.lastModifiedBy,
        "title": props.title, "subject": props.subject,
    }.items() if v}
    if verborgen:
        meta["_verborgen_tabbladen"] = ", ".join(verborgen)
    wb.close()
    return chunks, meta


def _pptx(path):
    from pptx import Presentation

    prs = Presentation(path)
    chunks = []
    for i, slide in enumerate(prs.slides, start=1):
        delen = []
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                delen.append(shape.text_frame.text)
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame.text.strip():
            delen.append("[notities] " + slide.notes_slide.notes_text_frame.text)
        if delen:
            chunks.append((f"dia {i}", "\n".join(delen)))
    cp = prs.core_properties
    meta = {k: v for k, v in {
        "author": cp.author, "last_modified_by": cp.last_modified_by,
        "title": cp.title, "subject": cp.subject,
    }.items() if v}
    return chunks, meta
